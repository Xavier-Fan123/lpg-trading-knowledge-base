"""Batch-extract text from pptx/ppt/xlsx/pdf files into 00_Inbox as raw .txt sources."""
import os, sys, hashlib, datetime

VAULT = r"C:\Users\itg\Desktop\lpg-trading-knowledge-base"
INBOX = os.path.join(VAULT, "00_Inbox")

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.replace("|", "").strip():
                        lines.append(row_text)
    return "\n".join(lines)

def extract_ppt_olefile(path):
    """Extract text from old .ppt format using olefile + raw byte scanning."""
    import olefile
    texts = []
    try:
        ole = olefile.OleFileIO(path)
        for stream_name in ole.listdir():
            joined = "/".join(stream_name)
            if "current user" in joined.lower():
                continue
            try:
                data = ole.openstream(stream_name).read()
                # Try UTF-16LE decoding for PowerPoint text streams
                if any(k in joined.lower() for k in ["document", "text", "powerpoint"]):
                    try:
                        decoded = data.decode("utf-16-le", errors="ignore")
                        # Filter printable content
                        cleaned = []
                        for ch in decoded:
                            if ch.isprintable() or ch in ("\n", "\r", "\t"):
                                cleaned.append(ch)
                            else:
                                cleaned.append(" ")
                        text = "".join(cleaned)
                        # Split into lines, keep meaningful ones
                        for line in text.split("\r"):
                            line = line.strip()
                            if len(line) > 2 and not all(c in " \t\x00" for c in line):
                                texts.append(line)
                    except:
                        pass
            except:
                pass
        ole.close()
    except Exception as e:
        texts.append(f"[OLE extraction partial - {e}]")

    if not texts:
        # Fallback: brute force scan for Chinese/ASCII text blocks
        with open(path, "rb") as f:
            raw = f.read()
        # Try to find UTF-16LE encoded Chinese text
        try:
            decoded = raw.decode("utf-16-le", errors="ignore")
            for line in decoded.split("\r"):
                line = line.strip()
                # Keep lines with Chinese chars or substantial ASCII
                has_cjk = any("\u4e00" <= c <= "\u9fff" for c in line)
                if (has_cjk and len(line) > 3) or (len(line) > 10 and line.isascii() and any(c.isalpha() for c in line)):
                    texts.append(line)
        except:
            pass

    return "\n".join(texts) if texts else "[Could not extract text from legacy .ppt format]"

def extract_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"\n=== Sheet: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            row_text = " | ".join(cells)
            if row_text.replace("|", "").strip():
                lines.append(row_text)
    return "\n".join(lines)

def extract_pdf(path):
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    lines = []
    total = len(reader.pages)
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            lines.append(f"\n--- Page {i+1}/{total} ---")
            lines.append(text.strip())
    return "\n".join(lines)

def save_to_inbox(filename, content):
    short_hash = hashlib.md5(filename.encode()).hexdigest()[:8]
    safe_name = os.path.splitext(filename)[0]
    # Keep Chinese chars but remove problematic filesystem chars
    for ch in ['/', '\\', ':', '*', '?', '"', '<', '>', '|']:
        safe_name = safe_name.replace(ch, '_')
    out_name = f"{safe_name}_{short_hash}.txt"
    out_path = os.path.join(INBOX, out_name)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(f"# Source: {filename}\n")
        f.write(f"# Extracted: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}\n")
        f.write(f"# Characters: {len(content)}\n\n")
        f.write(content)
    return out_name, len(content)

FILES = [
    (r"D:\Users\itg\xwechat_files\wxid_4360143601512_6947\msg\file\2026-04\船运执行.pptx", "pptx"),
    (r"D:\Users\itg\xwechat_files\wxid_4360143601512_6947\msg\file\2026-04\油轮合同.ppt", "ppt"),
    (r"D:\Users\itg\xwechat_files\wxid_4360143601512_6947\msg\file\2026-04\报关流程总结.ppt", "ppt"),
    (r"D:\Users\itg\xwechat_files\wxid_4360143601512_6947\msg\file\2026-04\BP,CHEVRON&SHELL GTC比较.xlsx", "xlsx"),
    (r"D:\Users\itg\xwechat_files\wxid_4360143601512_6947\msg\file\2026-04\GTC主要条款比较-Harold&Nadila.xlsx", "xlsx"),
    (r"D:\Users\itg\xwechat_files\wxid_4360143601512_6947\msg\file\2026-04\AWRP分享.pptx", "pptx"),
    (r"D:\Users\itg\xwechat_files\wxid_4360143601512_6947\msg\file\2026-04\石油贸易手册(1).pdf", "pdf"),
    (r"D:\Users\itg\xwechat_files\wxid_4360143601512_6947\msg\file\2026-04\石油加工基础知识070813.pdf", "pdf"),
]

extractors = {
    "pptx": extract_pptx,
    "ppt": extract_ppt_olefile,
    "xlsx": extract_xlsx,
    "pdf": extract_pdf,
}

print(f"Extracting {len(FILES)} files to {INBOX}\n")
for path, fmt in FILES:
    fname = os.path.basename(path)
    try:
        content = extractors[fmt](path)
        out_name, chars = save_to_inbox(fname, content)
        print(f"  OK  {fname} -> {out_name} ({chars:,} chars)")
    except Exception as e:
        print(f"  ERR {fname}: {e}")

print("\nDone.")
